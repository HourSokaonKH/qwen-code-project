from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create new presentation with 16:9 ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# NPIC Brand Colors
NPIC_BLUE = (0, 51, 102)      # #003366
NPIC_RED = (128, 0, 32)       # #800020
WHITE = (255, 255, 255)
LIGHT_GRAY = (240, 240, 240)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background rectangle (full blue)
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
    shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(*WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, image_path=None, image_right=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    
    # Content area
    if image_path and os.path.exists(image_path):
        if image_right:
            # Text on left
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6.5), Inches(5.5))
            # Image on right
            img = slide.shapes.add_picture(image_path, Inches(7.2), Inches(1.5), width=Inches(5.5))
        else:
            # Image on left
            img = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), width=Inches(5.5))
            # Text on right
            text_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(6.3), Inches(5.5))
        
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}" if isinstance(item, str) else item
            p.font.size = Pt(18)
            p.space_after = Pt(12)
            if isinstance(item, str):
                p.font.color.rgb = RGBColor(*NPIC_BLUE)
    else:
        # No image, full width text
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.333), Inches(5.5))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}" if isinstance(item, str) else item
            p.font.size = Pt(20)
            p.space_after = Pt(12)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    
    # Three columns for metrics
    col_width = Inches(4.0)
    start_x = Inches(0.7)
    
    for i, metric in enumerate(metrics):
        x_pos = start_x + (i * col_width)
        
        # Metric box
        metric_box = slide.shapes.add_shape(1, x_pos, Inches(1.5), col_width - Inches(0.2), Inches(4.5))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = RGBColor(*LIGHT_GRAY)
        metric_box.line.color.rgb = RGBColor(*NPIC_RED)
        metric_box.line.width = Pt(3)
        
        # Big number
        num_box = slide.shapes.add_textbox(x_pos + Inches(0.3), Inches(1.8), col_width - Inches(0.6), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric['number']
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*NPIC_RED)
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide.shapes.add_textbox(x_pos + Inches(0.3), Inches(3.3), col_width - Inches(0.6), Inches(2.0))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = metric['desc']
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(*NPIC_BLUE)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_comparison_slide(prs, title, comparison_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    
    # Table setup
    table_rows = len(comparison_data) + 1
    table_cols = 3
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.333)
    height = Inches(5.5)
    
    table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.4)
    table.columns[2].width = Inches(4.4)
    
    # Header row
    headers = ['Factor', 'Cloud Deployment', 'On-Premise']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(*WHITE)
            paragraph.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(comparison_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = RGBColor(*NPIC_BLUE)
                if col_idx == 0:
                    paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*LIGHT_GRAY)
    
    return slide

def add_roadmap_slide(prs, title, phases):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title bar
    title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.8))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    
    # Timeline arrow
    arrow = slide.shapes.add_shape(33, Inches(0.5), Inches(2.5), Inches(12.333), Inches(2.0))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(*LIGHT_GRAY)
    arrow.line.color.rgb = RGBColor(*NPIC_BLUE)
    arrow.line.width = Pt(2)
    
    # Phase markers
    phase_positions = [Inches(2.0), Inches(6.666), Inches(11.333)]
    
    for i, phase in enumerate(phases):
        x_pos = phase_positions[i]
        
        # Circle marker
        circle = slide.shapes.add_shape(9, x_pos - Inches(0.4), Inches(2.3), Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*NPIC_RED)
        circle.line.color.rgb = RGBColor(*NPIC_BLUE)
        circle.line.width = Pt(3)
        
        # Phase number
        num_box = slide.shapes.add_textbox(x_pos - Inches(0.3), Inches(2.4), Inches(0.6), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*WHITE)
        p.alignment = PP_ALIGN.CENTER
        
        # Phase details below
        detail_box = slide.shapes.add_textbox(x_pos - Inches(1.5), Inches(3.5), Inches(3.0), Inches(2.5))
        tf = detail_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{phase['title']}\n{phase['duration']}\n{phase['activities']}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*NPIC_BLUE)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_cta_slide(prs, title, cta_text, benefits):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*NPIC_BLUE)
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    # CTA button-like box
    cta_box = slide.shapes.add_shape(1, Inches(3.0), Inches(2.8), Inches(7.333), Inches(1.5))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = RGBColor(*NPIC_RED)
    cta_box.line.color.rgb = RGBColor(*WHITE)
    cta_box.line.width = Pt(2)
    
    cta_text_box = slide.shapes.add_textbox(Inches(3.2), Inches(3.1), Inches(6.933), Inches(0.9))
    tf = cta_text_box.text_frame
    p = tf.paragraphs[0]
    p.text = cta_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)
    p.alignment = PP_ALIGN.CENTER
    
    # Benefits list
    ben_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(2.0))
    tf = ben_box.text_frame
    tf.word_wrap = True
    
    for i, benefit in enumerate(benefits):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"✓ {benefit}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*WHITE)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
    
    return slide

# Build the presentation
print("Creating SMART-CAMPUS Executive Pitch...")

# Slide 1: Title
add_title_slide(prs, 
    "Transforming Education Through Technology",
    "SMART-CAMPUS: Enterprise Management System for NPIC")

# Slide 2: The Urgency
add_content_slide(prs,
    "The Cost of Inaction vs. Digital Transformation",
    [
        "Status Quo: Fragmented data, manual inefficiencies, revenue leakage",
        "Rising student expectations demand seamless digital experiences",
        "Competitive pressure: Modern institutions leverage technology",
        "The Future: Unified ecosystem, real-time insights, automated workflows"
    ],
    image_path="extracted_images/image_1.png",
    image_right=True)

# Slide 3: Solution Overview
add_content_slide(prs,
    "One Platform. Every Operation.",
    [
        "Unified Ecosystem: Admissions, Academics, Finance, HR in one system",
        "Role-Based Access: Granular permissions for all stakeholders",
        "Real-Time Analytics: Instant visibility into institutional performance",
        "Scalable Architecture: Grows with NPIC's ambitions"
    ],
    image_path="extracted_images/image_2.png",
    image_right=False)

# Slide 4: Key Benefits & ROI
metrics = [
    {'number': '85%', 'desc': 'Reduction in Administrative Processing Time'},
    {'number': '$0', 'desc': 'Revenue Shrinkage Through Automation'},
    {'number': '100%', 'desc': 'Data Integrity & Audit Trail'}
]
add_metrics_slide(prs, "Measurable Impact. Immediate ROI.", metrics)

# Slide 5: Academic Excellence
add_content_slide(prs,
    "Academic Excellence: Enrollment to Graduation",
    [
        "Digital admissions pipeline with document verification",
        "Automated course registration and scheduling",
        "Real-time grade management and transcript generation",
        "Comprehensive student lifecycle tracking"
    ],
    image_path="extracted_images/image_5.png",
    image_right=True)

# Slide 6: Facilities Management
add_content_slide(prs,
    "Facilities & Asset Management",
    [
        "Digital twin infrastructure monitoring",
        "IoT-enabled dormitory access control",
        "Preventive maintenance scheduling",
        "Asset tracking with QR code integration"
    ],
    image_path="extracted_images/image_10.png",
    image_right=False)

# Slide 7: Integrated Finance
add_content_slide(prs,
    "Integrated Financial Management",
    [
        "Auto-generated invoices with Bakong KHQR integration",
        "Real-time payment reconciliation",
        "Budget tracking and financial reporting",
        "Zero manual cash handling reduces risk"
    ],
    image_path="extracted_images/image_14.png",
    image_right=True)

# Slide 8: Cloud Deployment Strategy
comparison_data = [
    ['Upfront Costs', 'Low (Pay-as-you-go)', 'High (Hardware investment)'],
    ['Scalability', 'Instant elasticity', 'Fixed capacity'],
    ['Maintenance', 'Zero upkeep (managed)', '24/7 internal IT required'],
    ['Security', 'Built-in compliance', 'Your responsibility'],
    ['Uptime', '99.9%+ SLA guaranteed', 'Infrastructure dependent'],
    ['Updates', 'Automatic, seamless', 'Manual, risky']
]
add_comparison_slide(prs, "Cloud vs. On-Premise Deployment", comparison_data)

# Slide 9: Technology Foundation
add_content_slide(prs,
    "Enterprise-Grade Technology Stack",
    [
        "Next.js: Lightning-fast responsive interface",
        "NestJS: Robust, scalable backend architecture",
        "Docker Containers: Consistent deployment, 99.9% uptime",
        "MQTT Protocol: Real-time IoT device communication",
        "RBAC Security: Military-grade access control"
    ],
    image_path="extracted_images/image_18.png",
    image_right=False)

# Slide 10: Implementation Roadmap
phases = [
    {'title': 'Core Module Deployment', 'duration': 'Months 1-3', 'activities': 'Admissions, Registration, Finance'},
    {'title': 'Full Integration', 'duration': 'Months 4-6', 'activities': 'HR, Facilities, IoT sensors'},
    {'title': 'Go-Live & Optimization', 'duration': 'Months 7-9', 'activities': 'Training, rollout, refinement'}
]
add_roadmap_slide(prs, "Strategic 9-Month Rollout Plan", phases)

# Slide 11: Call to Action
benefits = [
    "Immediate 85% reduction in administrative overhead",
    "Zero revenue leakage through automated billing",
    "Future-ready cloud infrastructure",
    "Enhanced institutional reputation"
]
add_cta_slide(prs, 
    "Approve Phase 1 Implementation",
    "Authorize Cloud Deployment Initiation",
    benefits)

# Save the presentation
output_file = 'SMART_CAMPUS_Executive_Pitch_Final.pptx'
prs.save(output_file)
print(f"✓ Presentation saved: {output_file}")
print(f"✓ Total slides: {len(prs.slides)}")
print(f"✓ Slide size: 16:9 (13.333\" x 7.5\")")
print(f"✓ Images integrated from extracted_images folder")
